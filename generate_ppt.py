from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# Colors from template
DARK_BLUE = RGBColor(0x00, 0x33, 0x66)
GOLD = RGBColor(0xFF, 0xCC, 0x00)
RED = RGBColor(0xCC, 0x00, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
BLUE_HEADING = RGBColor(0x00, 0x33, 0x99)
LIGHT_BLUE = RGBColor(0x1F, 0x76, 0xB4)

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

DATE_TEXT = "5/16/2026"
INSTITUTION = "Dayananda Sagar Academy of Technology & Management\n(Autonomous Institute under VTU)"


def add_bottom_bar(slide):
    """Add the blue/yellow footer bar matching the template."""
    slide_w = prs.slide_width
    slide_h = prs.slide_height
    bar_h = Inches(0.45)

    # Dark blue left block (~35% width)
    left_box = slide.shapes.add_shape(
        1, 0, slide_h - bar_h, Inches(4.2), bar_h
    )
    left_box.fill.solid()
    left_box.fill.fore_color.rgb = RGBColor(0x00, 0x1F, 0x5B)
    left_box.line.fill.background()
    tf = left_box.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = DATE_TEXT
    run.font.color.rgb = RED
    run.font.size = Pt(9)
    run.font.bold = True

    # Gold/yellow right block
    right_box = slide.shapes.add_shape(
        1, Inches(4.2), slide_h - bar_h, slide_w - Inches(4.2), bar_h
    )
    right_box.fill.solid()
    right_box.fill.fore_color.rgb = GOLD
    right_box.line.fill.background()
    tf2 = right_box.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = "Dayananda Sagar Academy of Technology & Management\n(Autonomous Institute under VTU)"
    run2.font.color.rgb = BLACK
    run2.font.size = Pt(8)
    run2.font.bold = True


def add_page_number(slide, num):
    """Add centered page number in the gold bar area."""
    slide_h = prs.slide_height
    bar_h = Inches(0.45)
    txb = slide.shapes.add_textbox(Inches(6.3), slide_h - bar_h + Pt(2), Inches(0.5), bar_h)
    tf = txb.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = str(num)
    run.font.size = Pt(9)
    run.font.bold = True
    run.font.color.rgb = BLACK


def add_blue_line(slide, top_offset=Inches(1.15)):
    """Add the blue horizontal rule under the heading."""
    slide_w = prs.slide_width
    line = slide.shapes.add_shape(1, 0, top_offset, slide_w, Pt(3))
    line.fill.solid()
    line.fill.fore_color.rgb = DARK_BLUE
    line.line.fill.background()


def add_slide_heading(slide, text, top=Inches(0.7)):
    """Red bold heading at top-left."""
    txb = slide.shapes.add_textbox(Inches(0.2), top, Inches(10), Inches(0.5))
    tf = txb.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(24)
    run.font.bold = True
    run.font.color.rgb = RED
    return txb


def add_body_text(slide, text, left=Inches(0.4), top=Inches(1.3), width=Inches(12.5), height=Inches(5.5), size=Pt(16)):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf = txb.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    first = True
    for line in text.split("\n"):
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = line
        run.font.size = size
        run.font.color.rgb = BLACK
    return txb


def bullet_para(tf, text, bold_prefix=None, size=Pt(15), indent=False):
    p = tf.add_paragraph()
    p.space_before = Pt(6)
    if indent:
        p.level = 1
    if bold_prefix:
        r1 = p.add_run()
        r1.text = bold_prefix
        r1.font.bold = True
        r1.font.size = size
        r1.font.color.rgb = BLACK
        r2 = p.add_run()
        r2.text = text
        r2.font.size = size
        r2.font.color.rgb = BLACK
    else:
        r = p.add_run()
        r.text = text
        r.font.size = size
        r.font.color.rgb = BLACK
    return p


# ─────────────────────────────────────────────
# SLIDE 1 – TITLE
# ─────────────────────────────────────────────
slide1 = prs.slides.add_slide(prs.slide_layouts[6])

# Institution name top
top_txb = slide1.shapes.add_textbox(Inches(0.3), Inches(0.1), Inches(12.5), Inches(0.6))
tf = top_txb.text_frame
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "DAYANANDA SAGAR ACADEMY OF TECHNOLOGY & MANAGEMENT"
run.font.size = Pt(18)
run.font.bold = True
run.font.color.rgb = BLACK

p2 = tf.add_paragraph()
p2.alignment = PP_ALIGN.CENTER
r2 = p2.add_run()
r2.text = "(An Autonomous Institute under VTU)"
r2.font.size = Pt(13)
r2.font.color.rgb = BLACK

# Dept & semester
dept_txb = slide1.shapes.add_textbox(Inches(0.3), Inches(1.4), Inches(12.5), Inches(0.8))
tf2 = dept_txb.text_frame
p3 = tf2.paragraphs[0]
p3.alignment = PP_ALIGN.CENTER
r3 = p3.add_run()
r3.text = "Department of Information Science and Engineering"
r3.font.size = Pt(16)
r3.font.color.rgb = BLACK

p4 = tf2.add_paragraph()
p4.alignment = PP_ALIGN.CENTER
r4 = p4.add_run()
r4.text = "8th Semester-2026"
r4.font.size = Pt(15)
r4.font.color.rgb = BLACK

# Course code - purple
course_txb = slide1.shapes.add_textbox(Inches(0.3), Inches(2.25), Inches(12.5), Inches(0.4))
tf3 = course_txb.text_frame
p5 = tf3.paragraphs[0]
p5.alignment = PP_ALIGN.CENTER
r5 = p5.add_run()
r5.text = "Industry Internship [BINT803B]"
r5.font.size = Pt(16)
r5.font.bold = True
r5.font.color.rgb = RGBColor(0x7B, 0x2C, 0x9E)

# Horizontal rule
hr = slide1.shapes.add_shape(1, Inches(0.5), Inches(2.7), Inches(12.3), Pt(2))
hr.fill.solid()
hr.fill.fore_color.rgb = RGBColor(0xA0, 0x50, 0x00)
hr.line.fill.background()

# Title
title_txb = slide1.shapes.add_textbox(Inches(0.3), Inches(2.85), Inches(12.5), Inches(0.5))
tf4 = title_txb.text_frame
p6 = tf4.paragraphs[0]
p6.alignment = PP_ALIGN.CENTER
r6 = p6.add_run()
r6.text = "From Zero to Multi-Language Pro in 20 Hours"
r6.font.size = Pt(20)
r6.font.bold = True
r6.font.color.rgb = BLACK

# Sub-title / detail
sub_txb = slide1.shapes.add_textbox(Inches(0.3), Inches(3.35), Inches(12.5), Inches(0.4))
tf_sub = sub_txb.text_frame
p_sub = tf_sub.paragraphs[0]
p_sub.alignment = PP_ALIGN.CENTER
r_sub = p_sub.add_run()
r_sub.text = "The CodeXpert Approach – Bharat Unnati AI Fellowship"
r_sub.font.size = Pt(14)
r_sub.font.italic = True
r_sub.font.color.rgb = RGBColor(0x1F, 0x76, 0xB4)

# Name & USN
name_txb = slide1.shapes.add_textbox(Inches(0.3), Inches(3.85), Inches(12.5), Inches(0.7))
tf5 = name_txb.text_frame
p7 = tf5.paragraphs[0]
p7.alignment = PP_ALIGN.CENTER
r7 = p7.add_run()
r7.text = "Devarasetty Dinakar Sasank"
r7.font.size = Pt(16)
r7.font.color.rgb = LIGHT_BLUE

p8 = tf5.add_paragraph()
p8.alignment = PP_ALIGN.CENTER
r8 = p8.add_run()
r8.text = "[1DT22IS045]"
r8.font.size = Pt(15)
r8.font.color.rgb = LIGHT_BLUE

# Internal / External guide
int_txb = slide1.shapes.add_textbox(Inches(0.4), Inches(4.75), Inches(5), Inches(0.8))
tf6 = int_txb.text_frame
p9 = tf6.paragraphs[0]
r9 = p9.add_run()
r9.text = "Internal Guide -"
r9.font.size = Pt(12)
r9.font.color.rgb = BLACK
p10 = tf6.add_paragraph()
r10 = p10.add_run()
r10.text = "Asst. Professor"
r10.font.size = Pt(12)
r10.font.color.rgb = BLACK
p11 = tf6.add_paragraph()
r11 = p11.add_run()
r11.text = "Dept. of Information Science and Engineering"
r11.font.size = Pt(12)
r11.font.color.rgb = BLACK
p12 = tf6.add_paragraph()
r12 = p12.add_run()
r12.text = "DSATM"
r12.font.size = Pt(12)
r12.font.color.rgb = BLACK

ext_txb = slide1.shapes.add_textbox(Inches(8.5), Inches(4.75), Inches(4.5), Inches(0.4))
tf7 = ext_txb.text_frame
p13 = tf7.paragraphs[0]
p13.alignment = PP_ALIGN.LEFT
r13 = p13.add_run()
r13.text = "External Guide -"
r13.font.size = Pt(12)
r13.font.color.rgb = BLACK

add_bottom_bar(slide1)
add_page_number(slide1, 1)

# ─────────────────────────────────────────────
# SLIDE 2 – OUTLINE
# ─────────────────────────────────────────────
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_heading(slide2, "Outline")
add_blue_line(slide2)

txb = slide2.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(12.2), Inches(5.5))
tf = txb.text_frame
tf.word_wrap = True

items = [
    "1.  Introduction",
    "2.  Organization profile and background",
    "3.  Internship objectives and scope of work",
    "4.  Tools / Technologies / Methodologies Learnt",
    "5.  Weekly activity log / work description",
    "6.  Outcomes achieved and skills acquired",
    "7.  Challenges faced and solutions implemented",
    "8.  Results",
    "9.  Conclusion and future scope",
]
for i, item in enumerate(items):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_before = Pt(8)
    run = p.add_run()
    run.text = item
    run.font.size = Pt(17)
    run.font.color.rgb = BLACK

add_bottom_bar(slide2)
add_page_number(slide2, 2)

# ─────────────────────────────────────────────
# SLIDE 3 – INTRODUCTION
# ─────────────────────────────────────────────
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_heading(slide3, "Introduction")
add_blue_line(slide3)

txb = slide3.shapes.add_textbox(Inches(0.4), Inches(1.3), Inches(12.5), Inches(5.5))
tf = txb.text_frame
tf.word_wrap = True

intro_lines = [
    ("Background & Motivation", True),
    ("", False),
    ("India faces a growing skills gap with 40 million+ students entering the workforce annually, where over", False),
    ("85% lack programming confidence and 92% need AI literacy — yet quality training remains expensive", False),
    ("and inaccessible to most.", False),
    ("", False),
    ("The Bharat Unnati AI Fellowship, powered by the CodeXpert platform, was conceived as a national", False),
    ("upliftment mission to bridge this gap. It delivers a 20-hour intensive multi-language programming", False),
    ("course, 40-hour GenAI specialization, and 60-hour Agentic AI internship — all for just ₹499.", False),
    ("", False),
    ("The CodeXpert Approach follows a proven 10-step methodology:", False),
    ("Tokens → Statements → Practice → Tasks → OOP → Data Structures → Libraries →", False),
    ("Design Patterns → Vibe Coding → UI Development", False),
    ("", False),
    ("This internship provided hands-on exposure to 15+ programming languages, AI-assisted coding,", False),
    ("automation agent development, and real-world capstone project implementation.", False),
]

first = True
for text, bold in intro_lines:
    if first:
        p = tf.paragraphs[0]
        first = False
    else:
        p = tf.add_paragraph()
    run = p.add_run()
    run.text = text
    run.font.size = Pt(14)
    run.font.bold = bold
    run.font.color.rgb = DARK_BLUE if bold else BLACK

add_bottom_bar(slide3)
add_page_number(slide3, 3)

# ─────────────────────────────────────────────
# SLIDE 4 – ORGANIZATION PROFILE
# ─────────────────────────────────────────────
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_heading(slide4, "Organization Profile and Background")
add_blue_line(slide4)

txb = slide4.shapes.add_textbox(Inches(0.4), Inches(1.3), Inches(12.5), Inches(5.5))
tf = txb.text_frame
tf.word_wrap = True

org_content = [
    ("Organization:  ", "Learners Byte Pvt. Ltd. — Expertpedia AI Platform"),
    ("Recognition:  ", "AICTE/NEAT 6.0 Ecosystem — National Educational Alliance for Technology"),
    ("Mission:  ", "Democratizing AI and programming education for India's youth through affordable, structured, industry-aligned fellowships"),
    ("", ""),
    ("Key Offerings:", ""),
    ("", "• CodeXpert Track  — 20-hour multi-language programming intensive"),
    ("", "• Expertpedia AI  — Unlimited adaptive AI learning platform"),
    ("", "• GenAI Specialization  — 40-hour generative AI mastery program"),
    ("", "• Agentic AI Internship  — 60-hour automation & agent-building program"),
    ("", ""),
    ("Reach & Impact:", ""),
    ("", "• Targets 40M+ students entering the Indian workforce"),
    ("", "• Available to all branches: CSE, ECE, Mech, Civil, and more"),
    ("", "• 100% online with all session recordings provided"),
    ("", "• Verifiable certificate with unique QR code on completion"),
]

first = True
for bold_part, normal_part in org_content:
    if first:
        p = tf.paragraphs[0]
        first = False
    else:
        p = tf.add_paragraph()
    p.space_before = Pt(4)
    if bold_part:
        r1 = p.add_run()
        r1.text = bold_part
        r1.font.bold = True
        r1.font.size = Pt(14)
        r1.font.color.rgb = DARK_BLUE if ":" in bold_part else BLACK
    if normal_part:
        r2 = p.add_run()
        r2.text = normal_part
        r2.font.size = Pt(14)
        r2.font.color.rgb = BLACK

add_bottom_bar(slide4)
add_page_number(slide4, 4)

# ─────────────────────────────────────────────
# SLIDE 5 – OBJECTIVES & SCOPE
# ─────────────────────────────────────────────
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_heading(slide5, "Internship Objectives and Scope of Work")
add_blue_line(slide5)

txb = slide5.shapes.add_textbox(Inches(0.4), Inches(1.3), Inches(12.5), Inches(5.5))
tf = txb.text_frame
tf.word_wrap = True

obj_lines = [
    ("Objectives:", True, False),
    ("", False, False),
    ("Objective 1: ", "Achieve multi-language programming fluency across 15+ languages by mastering the 10-step CodeXpert methodology — covering tokens, syntax, OOP, data structures, libraries, and design patterns."),
    ("Objective 2: ", "Design and implement a real-world capstone project applying OOP principles, data structures, UI development, and design patterns in at least two different programming languages."),
    ("Objective 3: ", "Develop hands-on proficiency in AI-assisted coding (GenAI tools), build automation agents using CrewAI and LangChain, and create a portfolio of 13+ industry-relevant projects."),
    ("", ""),
    ("Scope of Work:", True, False),
    ("", ""),
    ("• Module 1–6: ", "20-hour CodeXpert Track — Tokens, Statements, Tasks, OOP, Libraries, Design Patterns, UI Development"),
    ("• Module 7: ", "40-hour GenAI Specialization — Advanced prompting, AI content creation, multimodal AI workflows"),
    ("• Module 8: ", "60-hour Agentic AI Internship — CrewAI, LangChain, 10+ automation portfolio projects"),
    ("• Capstone: ", "Full-stack project implementation with OOP, data structures, design patterns & UI in 2 languages"),
]

first = True
for item in obj_lines:
    if first:
        p = tf.paragraphs[0]
        first = False
    else:
        p = tf.add_paragraph()
    p.space_before = Pt(5)
    if len(item) == 3:
        text, bold, _ = item
        run = p.add_run()
        run.text = text
        run.font.bold = bold
        run.font.size = Pt(15)
        run.font.color.rgb = DARK_BLUE if bold else BLACK
    else:
        bold_t, normal_t = item
        if bold_t:
            r1 = p.add_run()
            r1.text = bold_t
            r1.font.bold = True
            r1.font.size = Pt(14)
            r1.font.color.rgb = DARK_BLUE
        if normal_t:
            r2 = p.add_run()
            r2.text = normal_t
            r2.font.size = Pt(14)
            r2.font.color.rgb = BLACK

add_bottom_bar(slide5)
add_page_number(slide5, 5)

# ─────────────────────────────────────────────
# SLIDE 6 – TOOLS / TECHNOLOGIES
# ─────────────────────────────────────────────
slide6 = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_heading(slide6, "Tools / Technologies / Methodologies Learnt")
add_blue_line(slide6)

txb = slide6.shapes.add_textbox(Inches(0.4), Inches(1.3), Inches(12.5), Inches(5.5))
tf = txb.text_frame
tf.word_wrap = True

tech_lines = [
    ("Programming Languages (15+):", True),
    ("Python, Java, C++, JavaScript, TypeScript, C#, Go, Rust, Kotlin, Swift, PHP, Ruby, Scala, Dart, R", False),
    ("", False),
    ("AI & Automation Frameworks:", True),
    ("CrewAI, LangChain, OpenAI API, Google Gemini API, Hugging Face Transformers, Anthropic Claude API", False),
    ("", False),
    ("Development Tools & Platforms:", True),
    ("VS Code, GitHub, Jupyter Notebooks, Replit, Google Colab, Expertpedia AI Platform", False),
    ("", False),
    ("Methodologies:", True),
    ("• 10-Step CodeXpert Methodology: Tokens → Statements → Practice → Tasks → OOP → Data Structures → Libraries → Design Patterns → Vibe Coding → UI Development", False),
    ("• Agile Sprint-based Learning with weekly deliverables", False),
    ("• Project-Based Learning (PBL) with capstone project", False),
    ("• AI-Assisted Development (Vibe Coding) using GenAI tools", False),
    ("", False),
    ("Key Concepts:", True),
    ("Object-Oriented Programming, Data Structures & Algorithms, REST APIs, Design Patterns (MVC, Singleton, Observer, Factory), UI/UX Fundamentals", False),
]

first = True
for text, bold in tech_lines:
    if first:
        p = tf.paragraphs[0]
        first = False
    else:
        p = tf.add_paragraph()
    p.space_before = Pt(4)
    run = p.add_run()
    run.text = text
    run.font.bold = bold
    run.font.size = Pt(13)
    run.font.color.rgb = DARK_BLUE if bold else BLACK

add_bottom_bar(slide6)
add_page_number(slide6, 6)

# ─────────────────────────────────────────────
# SLIDE 7 – WEEKLY LOG (Part 1)
# ─────────────────────────────────────────────
slide7 = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_heading(slide7, "Weekly Activity Log / Work Description")
add_blue_line(slide7)

txb = slide7.shapes.add_textbox(Inches(0.4), Inches(1.3), Inches(12.5), Inches(5.6))
tf = txb.text_frame
tf.word_wrap = True

weekly_1 = [
    ("Week 1 (Jan 12 – Jan 18)", ""),
    ("", "Orientation to the Bharat Unnati AI Fellowship; completed onboarding on the Expertpedia AI platform; introduced to the CodeXpert 10-step methodology and the program roadmap covering 15+ languages."),
    ("Week 2 (Jan 19 – Jan 25)", ""),
    ("", "Module 1 — Tokens & Syntax: Studied tokens, keywords, operators, and syntax rules across Python, Java, C++, and JavaScript; completed 5 syntax labs comparing constructs across languages."),
    ("Week 3 (Jan 26 – Feb 1)", ""),
    ("", "Module 2 — Statements & Multi-Language Practice: Built conditional, loop, and function statements in 6 languages; completed 4 cross-language translation exercises; submitted first mini-project (Student Performance Analyzer in Python)."),
    ("Week 4 (Feb 2 – Feb 8)", ""),
    ("", "Module 3 — Task Accomplishment: Applied learned syntax to solve 20+ algorithmic tasks; practiced debugging techniques; built a Daily Expense Tracker in JavaScript and an AI Language Converter in Python using the Gemini API."),
    ("Week 5 (Feb 9 – Feb 15)", ""),
    ("", "Module 4 — OOP Principles & Data Structures: Implemented classes, inheritance, polymorphism, and encapsulation in Python, Java, and C++; practiced stacks, queues, linked lists, trees, and hash maps with language-specific libraries."),
]

first = True
for bold_t, normal_t in weekly_1:
    if first:
        p = tf.paragraphs[0]
        first = False
    else:
        p = tf.add_paragraph()
    p.space_before = Pt(5)
    if bold_t:
        r = p.add_run()
        r.text = bold_t
        r.font.bold = True
        r.font.size = Pt(13)
        r.font.color.rgb = BLACK
    if normal_t:
        r2 = p.add_run()
        r2.text = normal_t
        r2.font.size = Pt(13)
        r2.font.color.rgb = BLACK

add_bottom_bar(slide7)
add_page_number(slide7, 7)

# ─────────────────────────────────────────────
# SLIDE 8 – WEEKLY LOG (Part 2)
# ─────────────────────────────────────────────
slide8 = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_heading(slide8, "Weekly Activity Log / Work Description")
add_blue_line(slide8)

txb = slide8.shapes.add_textbox(Inches(0.4), Inches(1.3), Inches(12.5), Inches(5.6))
tf = txb.text_frame
tf.word_wrap = True

weekly_2 = [
    ("Week 6 (Feb 16 – Feb 22)", ""),
    ("", "Module 5 — Libraries & Real-World Implementation: Integrated third-party libraries (NumPy, Pandas, React, Spring Boot); built an Invoice Automation Bot and an Order Management System using Python and JavaScript."),
    ("Week 7 (Feb 23 – Mar 1)", ""),
    ("", "Module 6 — Design Patterns, Vibe Coding & UI Development: Implemented MVC, Singleton, Observer, and Factory patterns; explored AI-assisted (Vibe Coding) development workflow; built a Task Management System UI in React with Python backend."),
    ("Week 8 (Mar 2 – Mar 8)", ""),
    ("", "GenAI Specialization: Completed advanced prompt engineering techniques; created AI-generated content, diagrams, and scripts; practiced multimodal AI tools and built productivity automation workflows."),
    ("Week 9 (Mar 9 – Mar 15)", ""),
    ("", "Agentic AI — Phase 1: Introduced to CrewAI and LangChain frameworks; built a Meeting Summarizer AI agent and a Task Planner Agent; designed multi-agent orchestration workflows."),
    ("Week 10 (Mar 16 – Mar 22)", ""),
    ("", "Agentic AI — Phase 2 & Capstone: Completed 10 automation portfolio projects; finalized the capstone Personal Finance Tracker with analytics dashboard, implemented in both Python (Tkinter) and JavaScript (React); obtained verifiable fellowship certificate."),
]

first = True
for bold_t, normal_t in weekly_2:
    if first:
        p = tf.paragraphs[0]
        first = False
    else:
        p = tf.add_paragraph()
    p.space_before = Pt(5)
    if bold_t:
        r = p.add_run()
        r.text = bold_t
        r.font.bold = True
        r.font.size = Pt(13)
        r.font.color.rgb = BLACK
    if normal_t:
        r2 = p.add_run()
        r2.text = normal_t
        r2.font.size = Pt(13)
        r2.font.color.rgb = BLACK

add_bottom_bar(slide8)
add_page_number(slide8, 8)

# ─────────────────────────────────────────────
# SLIDE 9 – OUTCOMES & SKILLS
# ─────────────────────────────────────────────
slide9 = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_heading(slide9, "Outcomes Achieved and Skills Acquired")
add_blue_line(slide9)

txb = slide9.shapes.add_textbox(Inches(0.4), Inches(1.3), Inches(12.5), Inches(5.5))
tf = txb.text_frame
tf.word_wrap = True

outcomes = [
    ("Technical Outcomes:", True),
    ("• Achieved programming fluency in 15+ languages — read, write, debug, and optimize code across any language", False),
    ("• Completed 20+ hands-on labs covering tokens, OOP, data structures, libraries, and design patterns", False),
    ("• Built 13+ portfolio projects including automation bots, a finance tracker, a task manager, and AI agents", False),
    ("• Implemented a full capstone project (Personal Finance Tracker) in Python and JavaScript with analytics UI", False),
    ("• Developed AI agents using CrewAI and LangChain capable of multi-step business task automation", False),
    ("", False),
    ("Skills Acquired:", True),
    ("• Multi-language programming proficiency — Python, Java, C++, JavaScript, TypeScript, Go, Rust, C# and more", False),
    ("• Object-Oriented Design, Data Structures & Algorithms, Design Patterns (MVC, Factory, Observer, Singleton)", False),
    ("• GenAI prompt engineering, AI content generation, and multimodal AI tool usage", False),
    ("• Agentic AI workflow design — building autonomous AI agents with tool use and memory", False),
    ("• UI Development fundamentals — building interactive interfaces with React and Tkinter", False),
    ("• Vibe Coding — AI-assisted software development methodology for rapid prototyping", False),
    ("• Team collaboration, version control (Git/GitHub), and professional project documentation", False),
]

first = True
for text, bold in outcomes:
    if first:
        p = tf.paragraphs[0]
        first = False
    else:
        p = tf.add_paragraph()
    p.space_before = Pt(4)
    run = p.add_run()
    run.text = text
    run.font.bold = bold
    run.font.size = Pt(13)
    run.font.color.rgb = DARK_BLUE if bold else BLACK

add_bottom_bar(slide9)
add_page_number(slide9, 9)

# ─────────────────────────────────────────────
# SLIDE 10 – CHALLENGES & SOLUTIONS
# ─────────────────────────────────────────────
slide10 = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_heading(slide10, "Challenges Faced and Solutions Implemented")
add_blue_line(slide10)

txb = slide10.shapes.add_textbox(Inches(0.4), Inches(1.3), Inches(12.5), Inches(5.5))
tf = txb.text_frame
tf.word_wrap = True

challenges = [
    ("Challenge 1: ", "Syntax overload while learning 15+ languages simultaneously\n"),
    ("Solution: ", "Adopted the CodeXpert token-first methodology — learning common patterns and mental models before language-specific syntax; used comparison charts to map constructs across languages.\n"),
    ("", ""),
    ("Challenge 2: ", "Understanding complex AI agent orchestration with CrewAI and LangChain\n"),
    ("Solution: ", "Started with single-agent tasks before progressing to multi-agent pipelines; used Expertpedia AI to generate explanatory diagrams and step-by-step breakdowns of framework internals.\n"),
    ("", ""),
    ("Challenge 3: ", "Integrating multiple technologies in the capstone project (backend + UI + database)\n"),
    ("Solution: ", "Applied the MVC design pattern to separate concerns; broke the project into sprints; used AI pair-programming (Vibe Coding) to accelerate boilerplate generation and focus on logic.\n"),
    ("", ""),
    ("Challenge 4: ", "Debugging cross-language behavior differences (e.g., Python dynamic typing vs. Java static typing)\n"),
    ("Solution: ", "Maintained a personal debug journal; leveraged AI tools for instant explanations of runtime errors; cross-referenced language documentation using Expertpedia AI's instant doubt resolution."),
]

first = True
for bold_t, normal_t in challenges:
    if first:
        p = tf.paragraphs[0]
        first = False
    else:
        p = tf.add_paragraph()
    p.space_before = Pt(4)
    if bold_t:
        r1 = p.add_run()
        r1.text = bold_t
        r1.font.bold = True
        r1.font.size = Pt(13)
        r1.font.color.rgb = DARK_BLUE
    if normal_t:
        r2 = p.add_run()
        r2.text = normal_t
        r2.font.size = Pt(13)
        r2.font.color.rgb = BLACK

add_bottom_bar(slide10)
add_page_number(slide10, 10)

# ─────────────────────────────────────────────
# SLIDE 11 – RESULTS
# ─────────────────────────────────────────────
slide11 = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_heading(slide11, "Results")
add_blue_line(slide11)

txb = slide11.shapes.add_textbox(Inches(0.4), Inches(1.3), Inches(12.5), Inches(5.5))
tf = txb.text_frame
tf.word_wrap = True

results = [
    ("Program Completion:", True),
    ("• Successfully completed all 6 CodeXpert modules (20 hours), GenAI Specialization (40 hours), and Agentic AI Internship (60 hours) — total 120 hours of structured learning", False),
    ("• Earned verifiable Bharat Unnati AI Fellowship Certificate with QR code and unique ID", False),
    ("", False),
    ("Portfolio Delivered:", True),
    ("• 3 Micro-Projects: Student Performance Analyzer, Daily Expense Tracker, AI Language Converter", False),
    ("• 10 Automation Projects: Invoice Automation Bot, Order Management System, Meeting Summarizer AI, Task Planner Agent, CRM Automation Tool, Document Q&A System, Email Workflow Automation, Data Collection Pipeline, Business Process Orchestrator, Multi-tool Integration Project", False),
    ("• Capstone: Personal Finance Tracker — expense tracking, analytics dashboard, report export (Python + JavaScript)", False),
    ("", False),
    ("Performance Metrics:", True),
    ("• 15+ programming languages mastered          • 20+ labs completed", False),
    ("• 13+ projects in portfolio                   • 4-month fellowship roadmap completed", False),
    ("• Capstone implemented in 2 languages          • AICTE/NEAT 6.0 aligned certification earned", False),
]

first = True
for text, bold in results:
    if first:
        p = tf.paragraphs[0]
        first = False
    else:
        p = tf.add_paragraph()
    p.space_before = Pt(5)
    run = p.add_run()
    run.text = text
    run.font.bold = bold
    run.font.size = Pt(13)
    run.font.color.rgb = DARK_BLUE if bold else BLACK

add_bottom_bar(slide11)
add_page_number(slide11, 11)

# ─────────────────────────────────────────────
# SLIDE 12 – CONCLUSION & FUTURE SCOPE
# ─────────────────────────────────────────────
slide12 = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_heading(slide12, "Conclusion and Future Scope")
add_blue_line(slide12)

txb = slide12.shapes.add_textbox(Inches(0.4), Inches(1.3), Inches(12.5), Inches(5.5))
tf = txb.text_frame
tf.word_wrap = True

conclusion = [
    ("Conclusion:", True),
    ("", False),
    ("The Bharat Unnati AI Fellowship provided a transformative learning experience that took me from limited programming knowledge to confident multi-language fluency and practical AI development skills.", False),
    ("", False),
    ("Through the CodeXpert methodology, I gained structured programming intuition applicable across any language — understanding that 90% of programming logic is universal, with only 10% being language-specific syntax.", False),
    ("", False),
    ("The hands-on agentic AI modules with CrewAI and LangChain equipped me with cutting-edge skills to build autonomous AI solutions relevant to real business problems — skills that place me ahead of the curve in India's rapidly evolving tech landscape.", False),
    ("", False),
    ("Future Scope:", True),
    ("", False),
    ("• Extend the capstone Personal Finance Tracker into a production-ready SaaS application with a cloud backend (AWS/GCP) and mobile UI (React Native / Flutter)", False),
    ("• Contribute to open-source AI agent projects on GitHub and participate in national hackathons", False),
    ("• Pursue advanced specializations in MLOps, LLM fine-tuning, and multi-agent system design", False),
    ("• Apply agentic AI skills to automate workflows in the ISE domain — intelligent code review bots, auto-grading agents, and research assistant AIs", False),
    ("• Mentor junior students in the CodeXpert methodology to scale the impact of the Bharat Unnati mission", False),
]

first = True
for text, bold in conclusion:
    if first:
        p = tf.paragraphs[0]
        first = False
    else:
        p = tf.add_paragraph()
    p.space_before = Pt(4)
    run = p.add_run()
    run.text = text
    run.font.bold = bold
    run.font.size = Pt(13)
    run.font.color.rgb = DARK_BLUE if bold else BLACK

add_bottom_bar(slide12)
add_page_number(slide12, 12)

# ─────────────────────────────────────────────
# SLIDE 13 – THANK YOU
# ─────────────────────────────────────────────
slide13 = prs.slides.add_slide(prs.slide_layouts[6])

txb = slide13.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.3), Inches(2))
tf = txb.text_frame
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "Thank You.."
run.font.size = Pt(60)
run.font.bold = True
run.font.color.rgb = LIGHT_BLUE

# bottom bar (no page number on thank-you per template style)
slide_w = prs.slide_width
slide_h = prs.slide_height
bar_h = Inches(0.45)

right_box = slide13.shapes.add_shape(1, 0, slide_h - bar_h, slide_w, bar_h)
right_box.fill.solid()
right_box.fill.fore_color.rgb = WHITE
right_box.line.fill.background()

inst_txb = slide13.shapes.add_textbox(Inches(5), slide_h - bar_h - Pt(5), Inches(8), Inches(0.5))
tf_i = inst_txb.text_frame
p_i = tf_i.paragraphs[0]
p_i.alignment = PP_ALIGN.RIGHT
r_i = p_i.add_run()
r_i.text = "Dayananda Sagar Academy of Technology & Management (Autonomous Institute under VTU)"
r_i.font.size = Pt(8)
r_i.font.bold = True
r_i.font.italic = True
r_i.font.color.rgb = BLACK

date_txb = slide13.shapes.add_textbox(Inches(0.1), slide_h - bar_h - Pt(5), Inches(2), Inches(0.4))
tf_d = date_txb.text_frame
p_d = tf_d.paragraphs[0]
r_d = p_d.add_run()
r_d.text = DATE_TEXT
r_d.font.size = Pt(9)
r_d.font.color.rgb = BLACK

num_txb = slide13.shapes.add_textbox(Inches(6.5), slide_h - bar_h - Pt(5), Inches(0.4), Inches(0.4))
tf_n = num_txb.text_frame
p_n = tf_n.paragraphs[0]
p_n.alignment = PP_ALIGN.CENTER
r_n = p_n.add_run()
r_n.text = "13"
r_n.font.size = Pt(9)
r_n.font.color.rgb = BLACK

# Save
output_path = r"c:\Users\palla\Desktop\SASANK INTERNSHIP\Sasank_Internship_Presentation.pptx"
prs.save(output_path)
print(f"Saved: {output_path}")
